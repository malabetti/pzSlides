import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from io import BytesIO
from flask import Response

def extrair_conteudo_wikipedia(topico):
    url = f"https://pt.wikipedia.org/wiki/{topico.replace(' ', '_')}"
    response = requests.get(url)
    soup = BeautifulSoup(response.text, 'html.parser')

    # Encontre o primeiro parágrafo de texto
    paragrafos = soup.find_all('p')

    if not paragrafos:
        return "ERRO"

    for paragrafo in paragrafos:
        p = paragrafo.get_text().strip()
        if len(p):
            return p
    return "ERRO"

def criar_slide(nome, titulo, subtitulo, num_slides, slide_content):
    doc = Presentation()

    slides = []

    slides.append(doc.slides.add_slide(doc.slide_layouts[0]))
    slides[0].placeholders[0].text = titulo
    slides[0].placeholders[1].text = subtitulo

    n = num_slides + 1

    number = 1
    for i in slide_content:
        t1 = i
        print(t1)
        slides.append(doc.slides.add_slide(doc.slide_layouts[1]))
        slides[number].placeholders[0].text = str(t1)

        t2 = slides[number].placeholders[1]

        res = extrair_conteudo_wikipedia(t1)
        #if res == "ERRO":
            #res = gpt(t1)

        t2.text = res
        number += 1

    #doc.save(f"{nome}.pptx")
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    response = Response(output.read())
    response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    response.headers['Content-Disposition'] = 'attachment; filename=apresentacao.pptx'

    return response
